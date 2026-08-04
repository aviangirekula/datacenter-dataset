"""Verify the poster layout and render a preview.

PowerPoint is not scriptable here, so the poster is checked directly from the
saved file: every shape's geometry is read back and tested against the rules
that actually get posters criticised (text off the page, boxes overlapping,
type below the legible floor, columns misaligned). A schematic preview is
rendered so the layout can also be seen.

    ./.venv/bin/python scripts/check_poster.py
"""
from __future__ import annotations

import sys
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.patches as mpatches  # noqa: E402
import matplotlib.pyplot as plt  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402

REPO = Path(__file__).resolve().parents[1]
PPTX = REPO / "figures" / "Oughton_Avilash_Angirekula_2026ASSIP_Poster.pptx"
PREVIEW = REPO / "figures" / "poster_preview.png"

W, H = 36.0, 27.0
# The template's own column geometry, read back from the shipped file. These are
# not ours to choose: the poster size and layout are fixed by the program.
COL_X = [0.41, 11.99, 24.13]
COL_R = [11.66, 23.78, 35.55]
CONTENT_BOTTOM = 26.54
MIN_PT = 19          # absolute floor anywhere on the poster
MIN_BODY_PT = 26     # program guidance; we target 28
HEADINGS = {"Background", "Materials and Methods", "Results", "Conclusions",
            "Major Citations", "Acknowledgements"}


def inches(v) -> float:
    return Emu(v).inches if v is not None else 0.0


# Deliberately measured a second time, more conservatively than the builder
# does, and with its own font metrics. If this shared the builder's estimator a
# bug in that estimator would hide itself, which is exactly what happened when
# the Conclusions bullet was clipped and this file reported no problems.
ARIAL = "/System/Library/Fonts/Supplemental/Arial.ttf"
ARIAL_BOLD = "/System/Library/Fonts/Supplemental/Arial Bold.ttf"
LINE_FACTOR = 1.20          # builder assumes 1.15; check against a stricter one
_F: dict = {}


def _measure(text: str, pt: float, bold: bool, width_in: float) -> int:
    from PIL import ImageFont
    key = (round(pt, 1), bold)
    if key not in _F:
        _F[key] = ImageFont.truetype(ARIAL_BOLD if bold else ARIAL,
                                     int(round(pt * 4)))
    f = _F[key]
    words = text.split()
    if not words:
        return 1
    limit = width_in * 72.0
    lines, cur = 1, ""
    for w in words:
        trial = w if not cur else f"{cur} {w}"
        if f.getlength(trial) / 4.0 <= limit:
            cur = trial
        else:
            lines, cur = lines + 1, w
    return lines


def overflow(shape) -> float:
    """Inches by which a shape's text exceeds its box. Negative means it fits."""
    tf = shape.text_frame
    if not tf.text.strip():
        return -99.0
    usable = inches(shape.width) - 0.20
    need = 0.10
    for p in tf.paragraphs:
        runs = [r for r in p.runs if r.text]
        if not runs:
            continue
        pt = max((r.font.size.pt for r in runs if r.font.size), default=18.0)
        bold = any(r.font.bold for r in runs)
        sp = p.line_spacing if isinstance(p.line_spacing, float) else 1.0
        n = _measure("".join(r.text for r in runs), pt, bold, usable)
        need += n * pt * LINE_FACTOR * sp / 72.0
        need += (p.space_after.pt if p.space_after else 0.0) / 72.0
    return need - inches(shape.height)


def main() -> None:
    prs = Presentation(str(PPTX))
    slide = prs.slides[0]
    problems, warnings = [], []

    if abs(prs.slide_width.inches - W) > 0.01 or abs(prs.slide_height.inches - H) > 0.01:
        problems.append(f"slide is {prs.slide_width.inches} x {prs.slide_height.inches}, must be 36 x 27")

    shapes = []
    for shp in slide.shapes:
        x, y = inches(shp.left), inches(shp.top)
        w, h = inches(shp.width), inches(shp.height)
        txt = ""
        sizes = []
        if shp.has_text_frame:
            txt = shp.text_frame.text
            for p in shp.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size:
                        sizes.append(r.font.size.pt)
        shapes.append({"x": x, "y": y, "w": w, "h": h, "txt": txt,
                       "sizes": sizes, "kind": str(shp.shape_type)})

        # 1. nothing may run off the page
        if x < -0.01 or y < -0.01 or x + w > W + 0.01 or y + h > H + 0.01:
            problems.append(
                f"off-page: '{txt[:34]}' at ({x:.2f},{y:.2f}) {w:.2f}x{h:.2f}")

        # 2. nothing may run past where the template's own boxes end
        if y > 4.7 and y + h > CONTENT_BOTTOM + 0.01 and txt.strip():
            problems.append(
                f"past the {CONTENT_BOTTOM} in content bottom: '{txt[:30]}' ends {y + h:.2f}")

        # 3. text must actually fit inside its own box
        if shp.has_text_frame:
            over = overflow(shp)
            if over > 0.02:
                problems.append(
                    f"text overflows its box by {over:.2f} in: '{txt[:34]}'")

        # 4. type floor
        for s in sizes:
            if s < MIN_PT:
                problems.append(f"{s:.0f}pt below the {MIN_PT}pt floor: '{txt[:34]}'")

    # 3. column alignment for body shapes
    for s in shapes:
        if s["y"] < 4.9 or not s["txt"].strip():
            continue                      # header band is its own geometry
        if s["w"] < 8.0:
            continue                      # sub-elements sit indented on purpose
        near_l = min(abs(s["x"] - c) for c in COL_X)
        if near_l > 0.30:
            warnings.append(
                f"left edge {s['x']:.2f} is off-grid: '{s['txt'][:30]}'")

    # 4. overlapping *filled* content boxes within a column
    def rect(s):
        return (s["x"], s["y"], s["x"] + s["w"], s["y"] + s["h"])

    body = [s for s in shapes if s["y"] > 4.9 and s["txt"].strip()]
    for i, a in enumerate(body):
        ax0, ay0, ax1, ay1 = rect(a)
        for b in body[i + 1:]:
            bx0, by0, bx1, by1 = rect(b)
            ox = min(ax1, bx1) - max(ax0, bx0)
            oy = min(ay1, by1) - max(ay0, by0)
            bar = a["txt"].strip() in HEADINGS or b["txt"].strip() in HEADINGS
            if ox > 0.5 and oy > (0.01 if bar else 0.25):
                (problems if bar else warnings).append(
                    f"overlap {ox:.2f}x{oy:.2f} in: '{a['txt'][:24]}' / '{b['txt'][:24]}'")

    # 5. word budget
    words = sum(len(s["txt"].split()) for s in shapes)
    if words > 620:
        warnings.append(f"{words} words total; the target is under ~600")

    # --- preview ---------------------------------------------------------------
    fig, ax = plt.subplots(figsize=(18, 13.5))
    ax.set_xlim(0, W)
    ax.set_ylim(H, 0)
    ax.set_aspect("equal")
    ax.axis("off")
    ax.add_patch(mpatches.Rectangle((0, 0), W, H, fc="white", ec="#999999", lw=2))
    for c in COL_X + COL_R:
        ax.axvline(c, color="#dddddd", lw=0.8, ls="--", zorder=0)
    for s in shapes:
        filled = not s["txt"].strip()
        ax.add_patch(mpatches.Rectangle(
            (s["x"], s["y"]), s["w"], s["h"],
            fc="#e8f0e8" if filled else "none",
            ec="#5a8f5a" if filled else "#bbbbbb", lw=1.0, zorder=1))
        if s["txt"].strip():
            size = max(s["sizes"]) if s["sizes"] else 20
            ax.text(s["x"] + 0.12, s["y"] + 0.30,
                    s["txt"][:70].replace("\n", " "),
                    fontsize=max(3.0, size * 0.115), color="#222222",
                    va="top", zorder=3)
    ax.set_title("Poster layout preview  (36 x 27 in, dashed = column grid)",
                 fontsize=13, color="#444444", pad=14)
    fig.savefig(PREVIEW, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)

    # --- report ----------------------------------------------------------------
    print(f"shapes: {len(shapes)} | words: {words}")
    print(f"smallest type: {min((min(s['sizes']) for s in shapes if s['sizes']), default=0):.0f}pt")
    print(f"preview: {PREVIEW.relative_to(REPO)}")
    if problems:
        print(f"\nPROBLEMS ({len(problems)}):")
        for p in problems:
            print("  x", p)
    if warnings:
        print(f"\nwarnings ({len(warnings)}):")
        for w in warnings[:14]:
            print("  !", w)
    if not problems and not warnings:
        print("\nAll layout checks passed.")
    sys.exit(1 if problems else 0)


if __name__ == "__main__":
    main()
