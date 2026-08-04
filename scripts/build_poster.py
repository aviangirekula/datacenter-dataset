"""Fill in the official ASSIP poster template.

This edits the program's own template in place. Every template shape keeps the
position, size, fill and heading style it ships with. Only the placeholder text
is replaced, box heights are adjusted to fit their content, and the figures are
placed in the space that frees up.

An earlier version of this script cleared the template's shapes and rebuilt a
similar-looking layout from scratch. The result looked close but was no longer
the institutional template, so this version touches nothing it does not have to.

Template geometry, read back from the shipped file and treated as fixed:
    col 1  x=0.41  w=11.25   Background (y=5.74), Materials and Methods (y=16.63)
    col 2  x=11.99 w=11.79   Results (y=5.74, one tall box)
    col 3  x=24.16 w=11.37   Conclusions (y=5.74), Citations (y=17.74),
                             Acknowledgements (y=23.66)

    ./.venv/bin/python scripts/build_poster.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[1]
TEMPLATE = Path("/tmp/assip_template.pptx")
FIGS = REPO / "figures" / "poster"
OUT = REPO / "figures" / "Oughton_Avilash_Angirekula_2026ASSIP_Poster.pptx"

GREEN = RGBColor(0x00, 0x66, 0x00)
GOLD = RGBColor(0xFF, 0xCC, 0x00)
LGREY = RGBColor(0xEF, 0xEF, 0xEF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DGREY = RGBColor(0x33, 0x33, 0x33)

BODY_PT = 26          # the template's own body size
BOTTOM = 26.54        # where the template's own content boxes end

TITLE = "A Multi-Hazard Assessment of 2,696 US Data Centers"
AUTHORS = "Avilash Angirekula¹, Dennies Bor¹, Edward J. Oughton¹"
AFFIL = ("¹Department of Geography and Geoinformation Sciences, "
         "College of Science, George Mason University")

BACKGROUND = (
    "•  AI and cloud demand has concentrated US computing into a few dozen "
    "metro clusters that face earthquakes, wildfire and flooding.\n"
    "•  Existing sources are proprietary, or give city-level coordinates never "
    "checked against a building.\n"
    "•  We built a free, public, building-level inventory of the contiguous US, "
    "and tested how much the answer moves when a coordinate is wrong.\n"
    "•  We measured exposure, meaning whether a site lies in a hazard zone. We "
    "did not model damage, so nothing here is a loss estimate."
)
BOUNDS = (
    "Limitations\n"
    "•  Coverage is uneven. 83% of our Virginia sites come from OpenStreetMap "
    "against 48% of our California ones, so part of the gap between states "
    "reflects who did the mapping.\n"
    "•  Hail and wind come from eyewitness reports, so within a state their "
    "rates track how built-up the surroundings are (rank correlation 0.39 and "
    "0.34). We left both out, along with lightning, tornado and hurricane.\n"
    "•  No public source lists power capacity, so a single server room counts "
    "the same as a 100 MW campus.\n"
    "•  40 sites shown as clear sit outside FEMA's mapped extent and 8 outside "
    "the wildfire raster, so 35% is a floor."
)
METHODS = (
    "We merged OpenStreetMap, PeeringDB and Wikidata into one record per site "
    "across the contiguous US, matching on name and distance, with written "
    "evidence kept for all 194 merges.\n"
    "Against FEMA USA Structures footprints, 1,681 coordinates fall inside a "
    "building, 993 match the nearest and 22 have no match. A second non-OSM "
    "source confirmed 1,546, and we re-checked the 1,150 least certain with an "
    "AI-assisted pipeline we then reviewed.\n"
    "Our first seismic layer failed validation against USGS reference values "
    "(22 of 150 sites within 10%), so we replaced it:"
)
METHODS_TABLE = [
    ("Earthquake", "USGS ASCE 7-22 shaking, 2% chance in 50 yr"),
    ("Wildfire", "USFS Hazard Potential 270 m, max in 2.4 km"),
    ("Flood", "FEMA NFHL layer 28, SFHA flag"),
    ("Water stress", "WRI Aqueduct 4.0 (a supply limit, reported apart)"),
    ("Uncertainty", "500 relocations per site, 10-500 m by precision"),
]
METHODS_TAIL = (
    "A site counts as exposed if it lies within 2.4 km of land rated High or "
    "Very High for wildfire, inside a FEMA Special Flood Hazard Area, or at "
    "peak ground acceleration of 0.30 g or above. Wildfire is buffered "
    "because 2,284 sites sit on land the wildfire map classes as developed, "
    "so the pixel under the building says nothing."
)

STATS = [("2,696", "facilities located"), ("954", "face a mapped hazard"),
         ("952", "in a high water-stress basin"), ("563", "water-stressed only")]
RESULTS_LEAD = ("The hazard map and the water map barely overlap. Counting "
                "water stress takes the fleet from 35% to 56% constrained.")
RESULTS_BODY = (
    "952 sites lie in a high or extremely-high water-stress basin, almost "
    "exactly the 954 facing a mapped hazard, but only 389 are the same sites. "
    "Water stress lands hardest where the hazard map says clear: Virginia is "
    "0.7% hazard-exposed and 31% water-stressed, Illinois 2.7% and 95%.\n"
    "Which hazard dominates depends on the region. New Jersey reaches 100% "
    "through wildfire alone and 0% through earthquake, California 96% "
    "earthquake. Nationally 646 sites are exposed to wildfire, 448 to "
    "earthquake and 71 to flood, and 205 face more than one.\n"
    "Widening the buffer to 5 km takes that count to 953, a 48% jump."
)
CONCLUSIONS = (
    "•  The 35% headline is the least useful number here. Only 7 of 32 states "
    "with 20 or more sites fall between 10% and 60% exposed.\n"
    "•  Water stress is not a natural hazard, but it constrains the same "
    "buildings, and it is high where the hazard map says clear.\n"
    "•  New Jersey and California are both 100% exposed for opposite reasons, "
    "so one national score would hide what a site faces.\n"
    "•  Next: fragility curves, to turn exposure into loss."
)
NULLS = (
    "Two checks on our choices\n"
    "Sampling across a building's whole footprint instead of its centre changed "
    "the answer for 29 of the 326 sites where both could be compared. Taller "
    "buildings looked more exposed nationally, but that vanished within states, "
    "so it was geography and not height."
)
CITATIONS = (
    "Dillon, G.K. (2023) Wildfire Hazard Potential for the United States, "
    "270-m, v2023. USDA Forest Service. doi:10.2737/RDS-2015-0047-4\n"
    "USGS (2026) ASCE 7-22 Seismic Design Maps web service. "
    "earthquake.usgs.gov/ws/designmaps\n"
    "FEMA (2026) National Flood Hazard Layer. hazards.fema.gov\n"
    "Oughton, E.J. and Weigel, R. (2026) A Comparative Multi-Hazard Risk "
    "Assessment of the US High-Voltage Transmission Network. Zenodo. "
    "doi:10.5281/zenodo.20331026 (CC BY 4.0)\n"
    "FEMA and ORNL (2024) USA Structures. gis.fema.gov\n"
    "WRI (2023) Aqueduct 4.0 Water Risk Atlas."
)
ACK = (
    "This research was made possible through the support of George Mason "
    "University's College of Science, which supports the ASSIP Program, and by "
    "the NCAR multi-hazards project."
)


ARIAL = "/System/Library/Fonts/Supplemental/Arial.ttf"
ARIAL_BOLD = "/System/Library/Fonts/Supplemental/Arial Bold.ttf"
# Arial's line box is taller than its point size. PowerPoint multiplies THIS by
# the line-spacing factor, not the point size, so leaving it out understates
# every box by about 15 per cent, which is what clipped the Conclusions bullet.
ARIAL_LINE = 1.15
TB_MARGIN_IN = 0.10          # text-frame left/right inset
_FONTS: dict = {}


def _font(size_pt: float, bold: bool):
    """Arial at 4x scale, so getlength() is precise to a quarter point."""
    from PIL import ImageFont
    key = (round(size_pt, 1), bold)
    if key not in _FONTS:
        _FONTS[key] = ImageFont.truetype(ARIAL_BOLD if bold else ARIAL,
                                         int(round(size_pt * 4)))
    return _FONTS[key]


def wrapped_lines(para: str, size_pt: float, width_in: float,
                  bold: bool = False) -> int:
    """How many lines this paragraph really takes, by measuring the glyphs."""
    words = para.split()
    if not words:
        return 1
    f = _font(size_pt, bold)
    limit_pt = width_in * 72.0
    lines, cur = 1, ""
    for w in words:
        trial = w if not cur else f"{cur} {w}"
        if f.getlength(trial) / 4.0 <= limit_pt:
            cur = trial
        else:
            lines += 1
            cur = w
    return lines


def fit_height(txt: str, size_pt: float, width_in: float,
               spacing: float = 0.95, space_after_pt: float = 10.0,
               bold: bool = False) -> float:
    """Height a block of text actually needs, in inches.

    Autofit is off, so PowerPoint does not shrink text that overflows, it clips
    it. Every line is measured with real Arial metrics rather than estimated
    from a character count, because the estimate ran about 15 per cent short and
    silently clipped text under the box below.
    """
    total = 0.0
    for para in txt.split("\n"):
        usable = width_in - 2 * TB_MARGIN_IN
        if para.lstrip().startswith("\u2022"):
            usable -= size_pt * 1.05 / 72.0     # hanging indent set by _bullet()
        n = wrapped_lines(para, size_pt, usable, bold=bold)
        total += n * size_pt * ARIAL_LINE * spacing / 72.0 + space_after_pt / 72.0
    return total + 2 * 0.05 + 0.18          # frame margins, then a little slack


def fit_size(txt: str, box_w: float, box_h: float, start_pt: float,
             bold: bool = True, spacing: float = 0.92) -> float:
    """Largest point size at which `txt` still fits the box, in 2 pt steps.

    The title bar is the template's, so the title adapts to the bar rather than
    the bar to the title.
    """
    pt = start_pt
    while pt > 24:
        if fit_height(txt, pt, box_w, spacing=spacing, bold=bold) <= box_h:
            return pt
        pt -= 2
    return 24.0


def fig_height(name: str, width_in: float) -> float:
    """Placed height from the image's real aspect ratio, never hardcoded."""
    from PIL import Image
    with Image.open(FIGS / f"{name}.png") as im:
        w, h = im.size
    return width_in * h / w


def find(slide, snippet: str):
    """The template shape whose placeholder text starts with `snippet`."""
    for shp in slide.shapes:
        if shp.has_text_frame and shp.text_frame.text.strip().startswith(snippet):
            return shp
    raise KeyError(f"template shape not found: {snippet!r}")


def geom(shape) -> tuple[float, float, float, float]:
    return (shape.left.inches, shape.top.inches,
            shape.width.inches, shape.height.inches)


def _bullet(p, line: str, size_pt: float) -> None:
    """Hang a bullet's wrapped lines under its text, not under the marker."""
    if not line.lstrip().startswith("\u2022"):
        return
    ind = int(round(size_pt * 12700 * 1.05))    # EMU per point is 12700
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(ind))
    pPr.set("indent", str(-ind))


def set_text(shape, txt: str, size_pt: float, bold=False, colour=BLACK,
             align=PP_ALIGN.LEFT, spacing=0.95) -> None:
    """Replace a template shape's text, leaving the shape itself alone."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(TB_MARGIN_IN)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    try:
        tf.auto_size = None            # never let PowerPoint silently shrink
    except Exception:  # noqa: BLE001 - older template parts lack the element
        pass
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    first = tf.paragraphs[0]
    for r in list(first.runs):
        r._r.getparent().remove(r._r)
    for i, line in enumerate(txt.split("\n")):
        p = first if i == 0 else tf.add_paragraph()
        p.alignment = align            # the template ships JUSTIFY, which
        p.line_spacing = spacing       # stretches word gaps on a wide column
        p.space_after = Pt(10)
        _bullet(p, line, size_pt)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = colour
        run.font.name = "Arial"


def new_text(slide, x, y, w, h, txt, size_pt, bold=False, colour=BLACK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=0.95,
             head_bold=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.10)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    for i, line in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(8)
        _bullet(p, line, size_pt)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size_pt + 1 if (head_bold and i == 0) else size_pt)
        r.font.bold = bold or (head_bold and i == 0)
        r.font.color.rgb = colour
        r.font.name = "Arial"
    return tb


def panel(slide, x, y, w, h, fill=LGREY, accent=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = RGBColor(0xA5, 0xA5, 0xA5)
    s.line.width = Pt(1)
    s.shadow.inherit = False
    if accent is not None:
        a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                   Inches(0.14), Inches(h))
        a.fill.solid()
        a.fill.fore_color.rgb = accent
        a.line.fill.background()
        a.shadow.inherit = False
    return s


def main() -> None:
    prs = Presentation(str(TEMPLATE))
    slide = prs.slides[0]
    assert abs(prs.slide_width.inches - 36) < 0.01, "template is not 36 in wide"
    assert abs(prs.slide_height.inches - 27) < 0.01, "template is not 27 in tall"

    # ---- header, edited in place ----
    ti = find(slide, "Put your title here")
    _, _, tw, th = geom(ti)
    tpt = fit_size(TITLE, tw, th, 80)
    set_text(ti, TITLE, tpt, bold=True, align=PP_ALIGN.CENTER, spacing=0.92)
    set_text(find(slide, "FirstName LastName1"), AUTHORS, 38, bold=True,
             align=PP_ALIGN.CENTER)
    set_text(find(slide, "1Department of Wherever"), AFFIL, 26,
             colour=DGREY, align=PP_ALIGN.CENTER)

    # ---- column 1: Background, then the limitations panel ----
    bg = find(slide, "Tell your audience")
    bx, by, bw, _ = geom(bg)
    bh = fit_height(BACKGROUND, BODY_PT, bw)
    set_text(bg, BACKGROUND, BODY_PT)
    bg.height = Inches(bh)

    nh = fit_height(BOUNDS, 23, bw - 0.45)
    # The Materials and Methods bar is fixed by the template at 15.72, so the
    # leftover space is split evenly rather than pooled into one dead gap.
    slack = (15.72 - 0.30) - (by + bh) - (nh + 0.16)
    py = by + bh + max(0.26, slack / 2)
    panel(slide, bx, py, bw, nh + 0.16, accent=GOLD)
    new_text(slide, bx + 0.26, py + 0.08, bw - 0.45, nh, BOUNDS, 22,
             head_bold=True)

    # ---- column 1: Materials and Methods ----
    me = find(slide, "Provide key details")
    mx, my, mw, _ = geom(me)
    mh = fit_height(METHODS, BODY_PT, mw)
    set_text(me, METHODS, BODY_PT)
    me.height = Inches(mh)

    ty = my + mh + 0.10
    for name, src in METHODS_TABLE:
        panel(slide, mx + 0.05, ty, mw - 0.10, 0.64)
        new_text(slide, mx + 0.20, ty + 0.01, 3.20, 0.62, name, 23, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
        new_text(slide, mx + 3.45, ty + 0.01, mw - 3.70, 0.62, src, 21,
                 colour=DGREY, anchor=MSO_ANCHOR.MIDDLE)
        ty += 0.69
    new_text(slide, mx, ty + 0.10, mw, fit_height(METHODS_TAIL, BODY_PT, mw),
             METHODS_TAIL, BODY_PT)

    # ---- column 2: Results. The template gives this one tall box, so the
    # figures are stacked inside its footprint. ----
    re_ = find(slide, "Describe the key data")
    rx, ry, rw, _ = geom(re_)
    y = ry

    sw = (rw - 0.30) / 4
    for i, (num, lab) in enumerate(STATS):
        x = rx + 0.05 + i * (sw + 0.08)
        new_text(slide, x, y, sw, 0.98, num, 48, bold=True,
                 colour=GREEN, align=PP_ALIGN.CENTER)
        new_text(slide, x, y + 0.92, sw, 0.86, lab, 20, colour=DGREY,
                 align=PP_ALIGN.CENTER)
    y += 1.86

    lh = fit_height(RESULTS_LEAD, 29, rw)
    set_text(re_, RESULTS_LEAD, 29, bold=True)
    re_.top, re_.height = Inches(y), Inches(lh)
    y += lh + 0.10

    slide.shapes.add_picture(str(FIGS / "fig1_map.png"), Inches(rx), Inches(y),
                             width=Inches(rw))
    y += fig_height("fig1_map", rw) + 0.08
    cap = ("Fig 1. One third of US data centers face a mapped hazard. Of the 205 "
           "facing two or more (orange), California holds 119, with smaller "
           "clusters in Utah, Nevada and New Jersey.")
    ch = fit_height(cap, 22, rw)
    new_text(slide, rx, y, rw, ch, cap, 22, bold=True)
    y += ch + 0.10

    rh = fit_height(RESULTS_BODY, BODY_PT, rw)
    new_text(slide, rx, y, rw, rh, RESULTS_BODY, BODY_PT)
    y += rh + 0.10

    slide.shapes.add_picture(str(FIGS / "fig2_states.png"), Inches(rx),
                             Inches(y), width=Inches(rw))
    y += fig_height("fig2_states", rw) + 0.06
    cap = "Fig 2. All 32 states with 20 or more sites, coloured by dominant hazard."
    new_text(slide, rx, y, rw, fit_height(cap, 22, rw), cap, 22, bold=True)

    # ---- column 3: Conclusions, then the null results and the QA figure ----
    co = find(slide, "Describe the key conclusions")
    cx, cy, cw, _ = geom(co)
    cch = fit_height(CONCLUSIONS, BODY_PT, cw)
    set_text(co, CONCLUSIONS, BODY_PT)
    co.height = Inches(cch)
    y = cy + cch + 0.20

    nh3 = fit_height(NULLS, 25, cw - 0.45)
    panel(slide, cx, y, cw, nh3 + 0.16, accent=GOLD)
    new_text(slide, cx + 0.26, y + 0.08, cw - 0.45, nh3, NULLS, 24,
             head_bold=True)
    y += nh3 + 0.30

    slide.shapes.add_picture(str(FIGS / "fig3_confidence.png"), Inches(cx),
                             Inches(y), width=Inches(cw))
    y += fig_height("fig3_confidence", cw) + 0.04
    cap = ("Fig 3. Moving each coordinate 500 times flips the wildfire class "
           "for at most 9% of sites placed within 100 m.")
    new_text(slide, cx, y, cw, fit_height(cap, 22, cw), cap, 22, bold=True)

    # ---- column 3: citations and acknowledgements, edited in place ----
    ci = find(slide, "Enter citations here")
    _, _, ciw, _ = geom(ci)
    set_text(ci, CITATIONS, 20, colour=DGREY, spacing=0.90)
    ci.height = Inches(fit_height(CITATIONS, 20, ciw, spacing=0.90))

    ak = find(slide, "This research was made possible")
    _, _, akw, _ = geom(ak)
    set_text(ak, ACK, 20, colour=DGREY)
    ak.height = Inches(fit_height(ACK, 20, akw))

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT.relative_to(REPO)}")
    print(f"  slide {prs.slide_width.inches:.0f} x {prs.slide_height.inches:.0f} in, "
          f"{len(slide.shapes)} shapes")


if __name__ == "__main__":
    main()
